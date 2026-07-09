#!/usr/bin/env python3
"""
Generador de presentación PPTX: Tutorial Docker y Kubernetes
Tema: Anthropic/Google (colores modernos, diseño limpio)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Colores tema Anthropic/Google
GOOGLE_BLUE = RGBColor(26, 115, 232)      # #1A73E8
GOOGLE_GREEN = RGBColor(15, 150, 110)     # #0F966E
ANTHROPIC_PURPLE = RGBColor(118, 75, 162) # #764BA2
LIGHT_GRAY = RGBColor(240, 244, 255)      # #F0F4FF
DARK_GRAY = RGBColor(60, 64, 67)          # #3C4043
WHITE = RGBColor(255, 255, 255)

def create_title_slide(prs, title, subtitle=""):
    """Crea diapositiva de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = GOOGLE_BLUE
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtítulo
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = LIGHT_GRAY
    
    return slide

def create_content_slide(prs, title, bullets, background_color=WHITE):
    """Crea diapositiva de contenido con bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = background_color
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = GOOGLE_BLUE
    
    # Línea divisora
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(0))
    line.line.color.rgb = GOOGLE_GREEN
    line.line.width = Pt(3)
    
    # Contenido
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

def create_code_slide(prs, title, code_text):
    """Crea diapositiva con código"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOOGLE_BLUE
    
    # Línea divisora
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(9), Inches(0))
    line.line.color.rgb = GOOGLE_GREEN
    line.line.width = Pt(3)
    
    # Código en caja oscura
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.3))
    code_frame = code_box.text_frame
    code_frame.word_wrap = True
    
    p = code_frame.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(11)
    p.font.name = 'Courier New'
    p.font.color.rgb = WHITE
    
    # Fondo oscuro para el código
    shape = slide.shapes.add_shape(1, Inches(0.35), Inches(1.15), Inches(9.3), Inches(5.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_GRAY
    shape.line.color.rgb = GOOGLE_GREEN
    shape.line.width = Pt(2)
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ===== SLIDE 1: Portada =====
    create_title_slide(prs, 
        "🐳 Docker & ☸️ Kubernetes", 
        "Tutorial completo paso a paso\nDesde instalación hasta orquestación en producción")
    
    # ===== SLIDE 2: Contenidos =====
    create_content_slide(prs, "Contenidos del tutorial", [
        "📋 Introducción: ¿Qué son Docker y Kubernetes?",
        "💾 Instalación en macOS, Linux y Windows",
        "🐳 Docker: Conceptos, imágenes, contenedores",
        "🔧 Construir y ejecutar tu primer contenedor",
        "☸️ Kubernetes: Orquestación de contenedores",
        "🚀 Desplegar aplicaciones en K8s",
        "📊 Monitoreo, escalado y mejores prácticas",
    ])
    
    # ===== SLIDE 3: ¿Qué es Docker? =====
    create_content_slide(prs, "¿Qué es Docker?", [
        "🐳 Plataforma de contenedorización",
        "📦 Empaqueta código + dependencias en una unidad",
        "🔄 \"Funciona en mi máquina\" → \"Funciona en todas\"",
        "⚡ Ligero, rápido y eficiente",
        "🌍 Permite ejecutar múltiples apps aisladas",
    ])
    
    # ===== SLIDE 4: ¿Qué es Kubernetes? =====
    create_content_slide(prs, "¿Qué es Kubernetes (K8s)?", [
        "☸️ Orquestador de contenedores",
        "🎯 Gestiona múltiples contenedores automáticamente",
        "⚖️ Distribuye carga, escala horizontalmente",
        "🔄 Auto-recuperación y actualizaciones sin downtime",
        "☁️ Estándar para infraestructura en la nube",
    ])
    
    # ===== SLIDE 5: Instalación macOS =====
    create_content_slide(prs, "Instalación en macOS", [
        "1️⃣ Opción A: Docker Desktop (Recomendado)",
        "   • Descarga: https://www.docker.com/products/docker-desktop",
        "   • Incluye Docker Engine + Kubernetes (opcional)",
        "",
        "2️⃣ Opción B: Homebrew",
        "   brew install docker docker-compose",
        "",
        "3️⃣ Verificar instalación:",
        "   docker --version",
        "   kubernetes version --client",
    ])
    
    # ===== SLIDE 6: Instalación Linux =====
    create_content_slide(prs, "Instalación en Linux (Ubuntu/Debian)", [
        "1️⃣ Actualizar paquetes:",
        "   sudo apt update && sudo apt upgrade",
        "",
        "2️⃣ Instalar Docker:",
        "   curl -fsSL https://get.docker.com -o get-docker.sh",
        "   sudo sh get-docker.sh",
        "",
        "3️⃣ Agregar usuario al grupo docker:",
        "   sudo usermod -aG docker $USER",
    ])
    
    # ===== SLIDE 7: Instalación Windows =====
    create_content_slide(prs, "Instalación en Windows", [
        "1️⃣ Requisitos: Windows 10/11 Pro, Enterprise o Education",
        "",
        "2️⃣ Habilitar Hyper-V:",
        "   • PowerShell (admin): Enable-WindowsOptionalFeature -FeatureName",
        "   • Hyper-V -All",
        "",
        "3️⃣ Instalar Docker Desktop:",
        "   https://www.docker.com/products/docker-desktop",
        "",
        "4️⃣ Reiniciar y abrir Docker Desktop",
    ])
    
    # ===== SLIDE 8: Conceptos Docker =====
    create_content_slide(prs, "Conceptos clave de Docker", [
        "🖼️ Imagen: Plantilla (sistema operativo + app)",
        "📦 Contenedor: Instancia en ejecución de una imagen",
        "🏗️ Dockerfile: Receta para construir una imagen",
        "🌐 Registry: Repositorio de imágenes (Docker Hub)",
        "📡 Volume: Almacenamiento persistente",
    ])
    
    # ===== SLIDE 9: Dockerfile =====
    create_code_slide(prs, "Ejemplo: Dockerfile", """FROM python:3.11-slim

WORKDIR /app

COPY app.py .
RUN pip install flask

EXPOSE 5000
ENV FLASK_APP=app.py

CMD ["python", "app.py"]""")
    
    # ===== SLIDE 10: Construir imagen Docker =====
    create_content_slide(prs, "Paso 1: Construir imagen Docker", [
        "📍 Ubicación: Terminal en raíz del proyecto",
        "",
        "1️⃣ Construir la imagen:",
        "   docker build -t mi-app-web:latest .",
        "",
        "✓ -t: Tag (nombre:versión)",
        "✓ .: Usar Dockerfile en directorio actual",
        "",
        "2️⃣ Verificar imagen creada:",
        "   docker images",
    ])
    
    # ===== SLIDE 11: Ejecutar contenedor =====
    create_content_slide(prs, "Paso 2: Ejecutar contenedor", [
        "🚀 Correr el contenedor en local:",
        "   docker run -p 8080:5000 mi-app-web:latest",
        "",
        "✓ -p 8080:5000: Mapea puerto local 8080 → puerto 5000",
        "",
        "🌐 Acceder en navegador:",
        "   http://localhost:8080",
        "",
        "✓ Deberías ver página web con interface moderna",
        "✓ Verifica status: \"✓ Contenedor ejecutándose correctamente\"",
    ])
    
    # ===== SLIDE 12: Verificar contenedor =====
    create_content_slide(prs, "Paso 3: Verificar contenedor", [
        "📊 Ver contenedores ejecutándose:",
        "   docker ps",
        "",
        "📝 Ver todos los contenedores (incluyendo detenidos):",
        "   docker ps -a",
        "",
        "🩺 Ver logs de la aplicación:",
        "   docker logs <container-id>",
        "",
        "⏹️ Detener contenedor:",
        "   docker stop <container-id>",
    ])
    
    # ===== SLIDE 13: Introducción K8s =====
    create_content_slide(prs, "Kubernetes: Conceptos clave", [
        "☸️ Cluster: Conjunto de máquinas trabajando juntas",
        "📌 Node: Máquina individual (física o virtual)",
        "📦 Pod: Unidad mínima (generalmente 1 contenedor)",
        "🎯 Deployment: Define cómo ejecutar pods",
        "📡 Service: Expone aplicaciones en la red",
        "⚙️ ConfigMap/Secret: Configuración e información sensible",
    ])
    
    # ===== SLIDE 14: Instalación Kubernetes =====
    create_content_slide(prs, "Instalación Kubernetes Local", [
        "🎯 Opción 1: Docker Desktop (más fácil)",
        "   • Settings → Kubernetes → Enable Kubernetes",
        "   • Esperar ~2 minutos",
        "",
        "🎯 Opción 2: Minikube (alternativa ligera)",
        "   brew install minikube",
        "   minikube start",
        "",
        "✅ Verificar instalación:",
        "   kubectl version --client",
        "   kubectl get nodes",
    ])
    
    # ===== SLIDE 15: Manifiestos K8s =====
    create_content_slide(prs, "Manifiestos Kubernetes", [
        "📄 YAML: Define infraestructura como código",
        "",
        "Tipos principales:",
        "  • Deployment: Ejecuta aplicaciones",
        "  • Service: Expone en la red",
        "  • ConfigMap: Datos de configuración",
        "  • PersistentVolume: Almacenamiento",
        "  • HorizontalPodAutoscaler: Auto-escalado",
    ])
    
    # ===== SLIDE 16: Desplegar en K8s =====
    create_content_slide(prs, "Paso 4: Desplegar en Kubernetes", [
        "📍 Requisitos previos:",
        "  1. Docker Desktop con K8s habilitado",
        "  2. Imagen Docker construida: mi-app-web:latest",
        "  3. Archivo kubernetes.yaml (proporcionado)",
        "",
        "🚀 Desplegar:",
        "   kubectl apply -f kubernetes.yaml",
        "",
        "✓ Crea Namespace, Deployment, Service y HPA",
    ])
    
    # ===== SLIDE 17: Verificar despliegue =====
    create_content_slide(prs, "Paso 5: Verificar despliegue en K8s", [
        "📊 Ver estado del despliegue:",
        "   kubectl get deployments -n docker-kubernetes-tutorial",
        "",
        "📦 Ver pods ejecutándose:",
        "   kubectl get pods -n docker-kubernetes-tutorial",
        "",
        "📡 Ver servicios:",
        "   kubectl get services -n docker-kubernetes-tutorial",
        "",
        "📝 Ver eventos/problemas:",
        "   kubectl describe pod <pod-name> -n docker-kubernetes-tutorial",
    ])
    
    # ===== SLIDE 18: Acceder a aplicación en K8s =====
    create_content_slide(prs, "Paso 6: Acceder a la aplicación", [
        "🔗 Port-forwarding (recomendado):",
        "   kubectl port-forward -n docker-kubernetes-tutorial \\",
        "     service/mi-app-web 8080:8080",
        "",
        "🌐 Acceder en navegador:",
        "   http://localhost:8080",
        "",
        "✓ Verás la misma página web que en Docker",
        "✓ Pero ahora ejecutándose en Kubernetes con 3 replicas",
    ])
    
    # ===== SLIDE 19: Escalar en K8s =====
    create_content_slide(prs, "Paso 7: Escalar aplicación", [
        "⚖️ Ver replicas actuales:",
        "   kubectl get deployments -n docker-kubernetes-tutorial",
        "",
        "📈 Escalar manualmente a 5 pods:",
        "   kubectl scale deployment/mi-app-web \\",
        "     --replicas=5 -n docker-kubernetes-tutorial",
        "",
        "📊 Auto-escalado (HPA):",
        "   ✓ Mínimo: 2 pods",
        "   ✓ Máximo: 5 pods",
        "   ✓ Trigger: CPU > 70%",
    ])
    
    # ===== SLIDE 20: Logs en K8s =====
    create_content_slide(prs, "Paso 8: Monitoreo y logs", [
        "📝 Ver logs de un pod:",
        "   kubectl logs <pod-name> -n docker-kubernetes-tutorial",
        "",
        "🔄 Seguir logs en tiempo real:",
        "   kubectl logs -f <pod-name> -n docker-kubernetes-tutorial",
        "",
        "🩺 Describir pod (información detallada):",
        "   kubectl describe pod <pod-name> \\",
        "     -n docker-kubernetes-tutorial",
    ])
    
    # ===== SLIDE 21: Limpiar recursos =====
    create_content_slide(prs, "Limpieza: Eliminar recursos", [
        "🗑️ Eliminar todo el despliegue en K8s:",
        "   kubectl delete namespace docker-kubernetes-tutorial",
        "",
        "🛑 Detener contenedor Docker:",
        "   docker stop <container-id>",
        "",
        "🗑️ Eliminar imagen Docker (opcional):",
        "   docker rmi mi-app-web:latest",
        "",
        "💾 Los datos están limpiados y listos para empezar",
    ])
    
    # ===== SLIDE 22: Mejores prácticas =====
    create_content_slide(prs, "✨ Mejores prácticas", [
        "🔒 Seguridad:",
        "  • No ejecutar como root en contenedores",
        "  • Usar secrets para datos sensibles",
        "  • Escanear imágenes para vulnerabilidades",
        "",
        "⚡ Rendimiento:",
        "  • Usar imágenes base ligeras (-slim, -alpine)",
        "  • Definir limits y requests",
        "  • Implementar health checks",
        "",
        "🔍 Observabilidad:",
        "  • Agregar logs estructurados",
        "  • Usar métricas y alertas",
        "  • Implementar trazas distribuidas",
    ])
    
    # ===== SLIDE 23: Compartir en GitHub =====
    create_content_slide(prs, "Compartir en GitHub (paso a paso)", [
        "1. Ver cambios:",
        "   git status",
        "",
        "2. Agregar archivos:",
        "   git add -A",
        "",
        "3. Crear commit:",
        "   git commit -m 'feat: agregar tutorial Docker y Kubernetes'",
        "",
        "4. Enviar a GitHub:",
        "   git push origin main",
        "",
        "5. Verificar: https://github.com/danielxxi/curso_git_diplomado",
    ])
    
    # ===== SLIDE 24: Troubleshooting =====
    create_content_slide(prs, "Troubleshooting - Solución de problemas", [
        "- Contenedor no inicia: docker logs <container-id>",
        "",
        "- Port ya en uso: docker run -p 9000:5000 ...",
        "",
        "- Pod en estado Pending: kubectl describe pod ...",
        "",
        "- K8s no funciona: kubectl cluster-info",
        "",
        "- Git error: Verifica rama correcta (git branch)",
    ])
    
    # ===== SLIDE 24: Recursos y referencias =====
    create_content_slide(prs, "📚 Recursos y referencias", [
        "🐳 Docker Official Docs:",
        "   https://docs.docker.com",
        "",
        "☸️ Kubernetes Official Docs:",
        "   https://kubernetes.io/docs",
        "",
        "📖 Docker Hub (imágenes):",
        "   https://hub.docker.com",
        "",
        "🛠️ Play with Kubernetes (online):",
        "   https://www.play-with-k8s.com",
    ])
    
    # ===== SLIDE 25: Resumen =====
    create_title_slide(prs,
        "🎉 ¡Completado!",
        "Ya sabes cómo:\n✓ Construir imágenes Docker\n✓ Ejecutar contenedores\n✓ Desplegar en Kubernetes\n✓ Escalar y monitorear aplicaciones")
    
    # Guardar presentación
    output_path = '/Users/daniel/curso_git_dpl/curso_git_diplomado/doc/Tutorial_Docker_Kubernetes.pptx'
    prs.save(output_path)
    print(f"✅ Presentación creada: {output_path}")

if __name__ == '__main__':
    main()
